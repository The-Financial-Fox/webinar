import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import os
from groq import Groq
from dotenv import load_dotenv
from operator import attrgetter
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Load API key securely
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

if not GROQ_API_KEY:
    st.error("🚨 API Key is missing! Set it in Streamlit Secrets or a .env file.")
    st.stop()

# Streamlit App UI
st.title("🤖 FP&A AI Agent - SaaS Cohort Analysis")
st.write("Upload an Excel file, analyze retention rates, and get AI-generated FP&A insights!")

# File uploader
uploaded_file = st.file_uploader("📂 Upload your cohort data (Excel format)", type=["xlsx"])

if uploaded_file:
    # Read the Excel file
    sales_data = pd.read_excel(uploaded_file)
    sales_data['Date'] = pd.to_datetime(sales_data['Date'])
    sales_data['CohortMonth'] = sales_data.groupby('Customer_ID')['Date'].transform('min').dt.to_period('M')
    sales_data['PurchaseMonth'] = sales_data['Date'].dt.to_period('M')
    sales_data['CohortIndex'] = (sales_data['PurchaseMonth'] - sales_data['CohortMonth']).apply(attrgetter('n'))
    
    # Create a pivot table for cohort analysis
    cohort_counts = sales_data.pivot_table(index='CohortMonth', columns='CohortIndex', values='Customer_ID', aggfunc='nunique')
    cohort_sizes = cohort_counts.iloc[:, 0]
    retention_rate = cohort_counts.divide(cohort_sizes, axis=0)
    
    # Identify outliers (clients who left all in one month)
    extreme_churn = sales_data.groupby(['PurchaseMonth'])['Customer_ID'].nunique()
    churn_outliers = extreme_churn[extreme_churn > (extreme_churn.mean() + 2 * extreme_churn.std())]
    
    # Plot retention rate heatmap and save to a buffer
    st.subheader("🔥 Retention Rate Heatmap")
    fig, ax = plt.subplots(figsize=(16, 9))
    sns.heatmap(retention_rate, annot=True, fmt=".0%", cmap="YlGnBu", linewidths=0.5, ax=ax)
    ax.set_title('Cohort Analysis - Retention Rate', fontsize=16)
    ax.set_xlabel('Months Since First Purchase', fontsize=12)
    ax.set_ylabel('Cohort Month', fontsize=12)
    plt.tight_layout()
    st.pyplot(fig)
    
    # Save the figure to a buffer
    heatmap_buffer = BytesIO()
    fig.savefig(heatmap_buffer, format='png', dpi=300)
    heatmap_buffer.seek(0)

    # AI Insights
    cohort_summary = f"""
    **Cohort Analysis Summary**:
    - Number of Cohorts: {len(cohort_counts)}
    - Retention Rate Breakdown:
    {retention_rate.to_string()}
    - Churn Outliers: {churn_outliers.to_string() if not churn_outliers.empty else 'No extreme churn detected.'}
    """
    
    st.subheader("🤖 AI Agent - FP&A Commentary")
    if st.button("🚀 Generate AI Commentary"):
        client = Groq(api_key=GROQ_API_KEY)
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are an AI-powered FP&A analyst providing deep financial insights, identifying key trends, anomalies, and actionable recommendations from cohort analysis data."
                },
                {
                    "role": "user",
                    "content": f"""
                        **Insight 1: Cohort Behavior & Retention Trends:** 
                        - Identify the distinct cohorts based on customer acquisition periods.
                        - Analyze retention trends across cohorts and highlight any unexpected patterns (e.g., sudden drops, consistent loyalty, or anomalies in renewal rates).

                        **Insight 2: Product-Level Retention & Revenue Impact:**  
                        - Determine which product has the highest and lowest retention rates and explain why.
                        - Evaluate the revenue impact of retention per product (i.e., do customers who retain longer contribute disproportionately to revenue?).
                        - Identify any upgrade/downgrade behaviors between products that affect retention and revenue.

                        **Insight 3: Extreme Churn & Outlier Detection:**  
                        - Detect customers with unusual churn patterns (e.g., large early drop-offs, reactivations, or one-time high spenders who disappear).
                        - Identify any customer segments (by product, cohort, or purchase behavior) that contribute to unexpected churn.
                        - Highlight any extreme cases where pricing anomalies (e.g., extreme discounts or bulk purchases) impact overall revenue patterns.

                        {cohort_summary}
                    """
                }
            ],
            model="llama3-8b-8192",
        )
        ai_commentary = response.choices[0].message.content
        
        # Display AI insights
        st.subheader("💡 AI-Generated FP&A Insights")
        st.write(ai_commentary)

        # PowerPoint Export Function
        def create_pptx(heatmap_img, commentary):
            prs = Presentation()

            # Slide 1: Title
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "SaaS Cohort Retention Analysis"
            subtitle.text = "Generated by FP&A AI Agent"

            # Slide 2: Retention Rate Heatmap
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Retention Rate Heatmap"

            # Add image
            img_path = "heatmap.png"
            with open(img_path, "wb") as f:
                f.write(heatmap_img.read())
            
            left = Inches(1)
            top = Inches(1.5)
            slide.shapes.add_picture(img_path, left, top, width=Inches(8))

            # Slide 3: AI Insights
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "AI-Generated FP&A Insights"

            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            text_frame.text = commentary

            # Save the PowerPoint file
            ppt_buffer = BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            return ppt_buffer

        # Generate PPT
        ppt_buffer = create_pptx(heatmap_buffer, ai_commentary)

        # Download button
        st.subheader("📤 Export Results")
        st.download_button(
            label="📊 Download PowerPoint Report",
            data=ppt_buffer,
            file_name="Cohort_Analysis_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
